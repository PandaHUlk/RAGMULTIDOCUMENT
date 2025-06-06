import streamlit as st
import os
import json
from datetime import datetime
from typing import List, Dict, Any
import pandas as pd
from PyPDF2 import PdfReader
# Set page config
st.set_page_config(
    page_title="Multi-Document RAG with FAQ System",
    page_icon="📚",
    layout="wide"
)

# Document processing imports

import docx
from pptx import Presentation
import csv
import openpyxl

# LangChain imports
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain.schema import Document

# Other imports
import google.generativeai as genai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

class DocumentProcessor:
    """Handles processing of various document types"""
    
    @staticmethod
    def extract_text_from_pdf(file) -> str:
        """Extract text from PDF files"""
        text = ""
        try:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        except Exception as e:
            st.error(f"Error reading PDF: {str(e)}")
        return text
    
    @staticmethod
    def extract_text_from_docx(file) -> str:
        """Extract text from DOCX files"""
        text = ""
        try:
            doc = docx.Document(file)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
        except Exception as e:
            st.error(f"Error reading DOCX: {str(e)}")
        return text
    
    @staticmethod
    def extract_text_from_pptx(file) -> str:
        """Extract text from PPTX files"""
        text = ""
        try:
            presentation = Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            st.error(f"Error reading PPTX: {str(e)}")
        return text
    
    @staticmethod
    def extract_text_from_txt(file) -> str:
        """Extract text from TXT files"""
        try:
            return file.read().decode('utf-8')
        except Exception as e:
            st.error(f"Error reading TXT: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_csv(file) -> str:
        """Extract text from CSV files"""
        text = ""
        try:
            df = pd.read_csv(file)
            text = df.to_string(index=False)
        except Exception as e:
            st.error(f"Error reading CSV: {str(e)}")
        return text
    
    @staticmethod
    def extract_text_from_excel(file) -> str:
        """Extract text from Excel files"""
        text = ""
        try:
            df = pd.read_excel(file, sheet_name=None)  # Read all sheets
            for sheet_name, sheet_df in df.items():
                text += f"Sheet: {sheet_name}\n"
                text += sheet_df.to_string(index=False) + "\n\n"
        except Exception as e:
            st.error(f"Error reading Excel: {str(e)}")
        return text

class FAQManager:
    """Handles FAQ functionality"""
    
    def __init__(self, faq_file="faqs.json"):
        self.faq_file = faq_file
        self.faqs = self.load_faqs()
    
    def load_faqs(self) -> List[Dict]:
        """Load FAQs from JSON file"""
        if os.path.exists(self.faq_file):
            try:
                with open(self.faq_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except:
                return []
        return []
    
    def save_faqs(self):
        """Save FAQs to JSON file"""
        with open(self.faq_file, 'w', encoding='utf-8') as f:
            json.dump(self.faqs, f, indent=2, ensure_ascii=False)
    
    def add_faq(self, question: str, answer: str, category: str = "General"):
        """Add a new FAQ"""
        new_faq = {
            "id": len(self.faqs) + 1,
            "question": question,
            "answer": answer,
            "category": category,
            "created_at": datetime.now().isoformat(),
            "usage_count": 0
        }
        self.faqs.append(new_faq)
        self.save_faqs()
    
    def search_faqs(self, query: str) -> List[Dict]:
        """Search FAQs based on query"""
        results = []
        query_lower = query.lower()
        for faq in self.faqs:
            if (query_lower in faq["question"].lower() or 
                query_lower in faq["answer"].lower()):
                results.append(faq)
        return results
    
    def get_faqs_by_category(self, category: str) -> List[Dict]:
        """Get FAQs by category"""
        return [faq for faq in self.faqs if faq["category"] == category]
    
    def increment_usage(self, faq_id: int):
        """Increment usage count for an FAQ"""
        for faq in self.faqs:
            if faq["id"] == faq_id:
                faq["usage_count"] += 1
                self.save_faqs()
                break

class RAGSystem:
    """Main RAG system with document processing and FAQ integration"""
    
    def __init__(self):
        self.doc_processor = DocumentProcessor()
        self.faq_manager = FAQManager()
        self.vector_store = None
        self.embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    def process_documents(self, uploaded_files) -> str:
        """Process multiple documents of various types"""
        all_text = ""
        processed_files = []
        
        for file in uploaded_files:
            file_type = file.type
            file_name = file.name
            
            try:
                if file_type == "application/pdf":
                    text = self.doc_processor.extract_text_from_pdf(file)
                elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    text = self.doc_processor.extract_text_from_docx(file)
                elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    text = self.doc_processor.extract_text_from_pptx(file)
                elif file_type == "text/plain":
                    text = self.doc_processor.extract_text_from_txt(file)
                elif file_type == "text/csv":
                    text = self.doc_processor.extract_text_from_csv(file)
                elif file_type in ["application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
                    text = self.doc_processor.extract_text_from_excel(file)
                else:
                    st.warning(f"Unsupported file type: {file_type} for {file_name}")
                    continue
                
                all_text += f"\n--- Content from {file_name} ---\n{text}\n"
                processed_files.append(file_name)
                
            except Exception as e:
                st.error(f"Error processing {file_name}: {str(e)}")
        
        st.success(f"Successfully processed {len(processed_files)} files: {', '.join(processed_files)}")
        return all_text
    
    def create_text_chunks(self, text: str) -> List[str]:
        """Split text into chunks"""
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        return text_splitter.split_text(text)
    
    def create_vector_store(self, text_chunks: List[str]):
        """Create and save vector store"""
        try:
            self.vector_store = FAISS.from_texts(text_chunks, self.embeddings)
            self.vector_store.save_local("faiss_index")
            st.success("Vector store created and saved successfully!")
        except Exception as e:
            st.error(f"Error creating vector store: {str(e)}")
    
    def load_vector_store(self):
        """Load existing vector store"""
        try:
            if os.path.exists("faiss_index"):
                self.vector_store = FAISS.load_local(
                    "faiss_index", 
                    self.embeddings, 
                    allow_dangerous_deserialization=True
                )
                return True
        except Exception as e:
            st.error(f"Error loading vector store: {str(e)}")
        return False
    
    def get_conversational_chain(self):
        """Create conversational chain"""
        prompt_template = """
        You are a helpful AI assistant. Answer the question based on the provided context.
        If the answer is not available in the context, clearly state that the information is not available in the provided documents.
        Provide detailed and accurate answers when possible.
        
        Context:
        {context}
        
        Question: {question}
        
        Answer:
        """
        
        model = ChatGoogleGenerativeAI(
            model="gemini-2.0-flash",
            temperature=0.3,
        )
        
        prompt = PromptTemplate(
            template=prompt_template, 
            input_variables=["context", "question"]
        )
        
        return load_qa_chain(
            llm=model,
            chain_type="stuff",
            prompt=prompt,
        )
    
    def process_query(self, question: str) -> str:
        """Process user query with FAQ integration"""
        # First, check FAQs
        faq_results = self.faq_manager.search_faqs(question)
        
        if faq_results:
            st.info("Found relevant FAQ(s):")
            for faq in faq_results[:3]:  # Show top 3 FAQ results
                with st.expander(f"FAQ: {faq['question']}"):
                    st.write(faq['answer'])
                    if st.button(f"Mark as helpful", key=f"faq_{faq['id']}"):
                        self.faq_manager.increment_usage(faq['id'])
                        st.success("Thank you for the feedback!")
        
        # Then search in documents
        if not self.vector_store:
            if not self.load_vector_store():
                return "Please upload and process documents first."
        
        try:
            docs = self.vector_store.similarity_search(question, k=4)
            chain = self.get_conversational_chain()
            
            response = chain(
                {"input_documents": docs, "question": question},
                return_only_outputs=True
            )
            
            return response["output_text"]
        except Exception as e:
            return f"Error processing query: {str(e)}"

def main():
    st.title("📚 Chat with SOP based SR")
    st.markdown("Upload documents of various formats and ask questions. The system also includes FAQ functionality.")
    
    # Initialize RAG system
    if 'rag_system' not in st.session_state:
        st.session_state.rag_system = RAGSystem()
    
    rag_system = st.session_state.rag_system
    
    # Sidebar for document upload and FAQ management
    with st.sidebar:
        st.header("📁 Document Management")
        
        # Document upload
        uploaded_files = st.file_uploader(
            "Upload Documents",
            type=["pdf", "docx", "pptx", "txt", "csv", "xlsx", "xls"],
            accept_multiple_files=True,
            help="Supported formats: PDF, DOCX, PPTX, TXT, CSV, XLSX"
        )
        
        if uploaded_files and st.button("🔄 Process Documents"):
            with st.spinner("Processing documents..."):
                raw_text = rag_system.process_documents(uploaded_files)
                if raw_text.strip():
                    text_chunks = rag_system.create_text_chunks(raw_text)
                    rag_system.create_vector_store(text_chunks)
                    st.balloons()
        
        st.divider()
        
        # FAQ Management
        st.header("❓ FAQ Management")
        
        with st.expander("Add New FAQ"):
            faq_question = st.text_input("Question")
            faq_answer = st.text_area("Answer")
            faq_category = st.selectbox(
                "Category",
                ["General", "Technical", "Process", "Policy", "Other"]
            )
            
            if st.button("Add FAQ") and faq_question and faq_answer:
                rag_system.faq_manager.add_faq(faq_question, faq_answer, faq_category)
                st.success("FAQ added successfully!")
        
        # FAQ Statistics
        if rag_system.faq_manager.faqs:
            st.metric("Total FAQs", len(rag_system.faq_manager.faqs))
    
    # Main chat interface
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.header("💬 Ask Questions")
        
        user_question = st.text_input(
            "Ask a question about your documents or general queries:",
            placeholder="e.g., What is the main topic discussed in the documents?"
        )
        
        if user_question:
            with st.spinner("Searching for answers..."):
                response = rag_system.process_query(user_question)
                
                st.subheader("📝 Answer from Documents:")
                st.write(response)
                
                # Option to add this Q&A as FAQ
                if st.button("💾 Save this as FAQ"):
                    rag_system.faq_manager.add_faq(
                        user_question, 
                        response, 
                        "Generated"
                    )
                    st.success("Added to FAQ database!")
    
    with col2:
        st.header("🔍 FAQ Browser")
        
        # Category filter
        categories = list(set([faq["category"] for faq in rag_system.faq_manager.faqs]))
        selected_category = st.selectbox("Filter by category:", ["All"] + categories)
        
        # Display FAQs
        faqs_to_show = (rag_system.faq_manager.faqs if selected_category == "All" 
                       else rag_system.faq_manager.get_faqs_by_category(selected_category))
        
        for faq in faqs_to_show[:5]:  # Show top 5
            with st.expander(f"❓ {faq['question'][:50]}..."):
                st.write(f"**Category:** {faq['category']}")
                st.write(f"**Answer:** {faq['answer']}")
                st.write(f"**Used:** {faq['usage_count']} times")

if __name__ == "__main__":
    main()
